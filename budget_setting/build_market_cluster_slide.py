from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from matplotlib.lines import Line2D
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
DATA_CSV = BASE_DIR / "market_cluster_aggressiveness_output.csv"
SCATTER_PNG = BASE_DIR / "market_clusters_scatter_static.png"
OUTPUT_PPTX = BASE_DIR / "market_cluster_aggressiveness_slide.pptx"


def load_data() -> pd.DataFrame:
    df = pd.read_csv(DATA_CSV)
    return df


def build_scatter_image(df: pd.DataFrame) -> None:
    plot_df = df[df["aggressiveness_uplift"].isin(["+60%", "+50%", "+40%"])].copy()
    colors = {"+60%": "#1B9E77", "+50%": "#D95F02", "+40%": "#7570B3"}

    fig, ax = plt.subplots(figsize=(10.8, 5.6), dpi=200)
    ax.set_facecolor("#F8FAFC")
    fig.patch.set_facecolor("#FFFFFF")

    for uplift, sub in plot_df.groupby("aggressiveness_uplift"):
        sizes = (sub["Budget_2026_EUR"] / sub["Budget_2026_EUR"].max()) * 1200 + 120
        ax.scatter(
            sub["market_gap_pct"],
            sub["opportunity_pct"],
            s=sizes,
            c=colors[uplift],
            alpha=0.9,
            edgecolors="#FFFFFF",
            linewidths=1.3,
            label=uplift,
        )
        for _, row in sub.iterrows():
            ax.text(
                row["market_gap_pct"] + 1.2,
                row["opportunity_pct"] + 0.8,
                row["Market"],
                fontsize=8.5,
                color="#111827",
            )

    x_med = plot_df["market_gap_pct"].median()
    y_med = plot_df["opportunity_pct"].median()
    ax.axvline(x=x_med, color="#64748B", linestyle="--", linewidth=1.0)
    ax.axhline(y=y_med, color="#64748B", linestyle="--", linewidth=1.0)

    ax.set_xlim(0, 105)
    ax.set_ylim(0, 105)
    ax.set_xlabel("Market Gap Percentile (higher = more under-allocated vs 2025 registrations)", fontsize=9)
    ax.set_ylabel("Opportunity Score Percentile (higher = stronger scaling potential)", fontsize=9)
    ax.set_title("Market Clusters: Opportunity Score vs Market Gap", fontsize=13, fontweight="bold", pad=10)
    ax.grid(alpha=0.18)
    # Use fixed-size custom legend markers to avoid clutter from large bubble sizes.
    legend_handles = [
        Line2D(
            [0],
            [0],
            marker="o",
            linestyle="None",
            markerfacecolor=colors["+40%"],
            markeredgecolor="#FFFFFF",
            markeredgewidth=0.8,
            markersize=9,
            label="+40%",
        ),
        Line2D(
            [0],
            [0],
            marker="o",
            linestyle="None",
            markerfacecolor=colors["+50%"],
            markeredgecolor="#FFFFFF",
            markeredgewidth=0.8,
            markersize=9,
            label="+50%",
        ),
        Line2D(
            [0],
            [0],
            marker="o",
            linestyle="None",
            markerfacecolor=colors["+60%"],
            markeredgecolor="#FFFFFF",
            markeredgewidth=0.8,
            markersize=9,
            label="+60%",
        ),
    ]
    ax.legend(
        handles=legend_handles,
        title="Uplift band",
        frameon=False,
        loc="lower right",
        fontsize=8.5,
        title_fontsize=9,
    )
    for spine in ax.spines.values():
        spine.set_color("#CBD5E1")
    plt.tight_layout()
    fig.savefig(SCATTER_PNG, bbox_inches="tight")
    plt.close(fig)


def _add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.13), Inches(12.2), Inches(0.56))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(17, 24, 39)


def _add_subtitle(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(12.2), Inches(0.42))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12.5)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(55, 65, 81)


def _add_method_box(slide) -> None:
    shape = slide.shapes.add_shape(
        autoshape_type_id=1, left=Inches(0.45), top=Inches(6.15), width=Inches(12.2), height=Inches(0.86)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(239, 246, 255)
    shape.line.color.rgb = RGBColor(209, 213, 219)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Method used: Markets are clustered on two dimensions — target pressure (DCFS gap vs 2026 target) and scalability "
        "(curve-based incremental response). Uplift bands (+40%, +50%, +60%) are directional aggressiveness signals."
    )
    p.font.size = Pt(10.5)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(31, 41, 55)


def _add_cluster_summary(slide, df: pd.DataFrame) -> None:
    summary = (
        df[df["aggressiveness_uplift"].isin(["+60%", "+50%", "+40%"])]
        .groupby("aggressiveness_uplift", dropna=False)
        .agg(markets=("Market", "count"), budget=("Budget_2026_EUR", "sum"))
        .reset_index()
        .sort_values("aggressiveness_uplift", ascending=False)
    )

    box = slide.shapes.add_textbox(Inches(10.05), Inches(1.18), Inches(2.9), Inches(0.5))
    t = box.text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = "Budget guidance bands"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(17, 24, 39)

    y = 1.58
    color_map = {"+60%": RGBColor(27, 158, 119), "+50%": RGBColor(217, 95, 2), "+40%": RGBColor(117, 112, 179)}
    for _, r in summary.iterrows():
        shape = slide.shapes.add_shape(1, Inches(10.05), Inches(y), Inches(2.9), Inches(0.62))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(229, 231, 235)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{r['aggressiveness_uplift']}  |  {int(r['markets'])} markets  |  €{int(r['budget']):,}"
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.color.rgb = color_map[str(r["aggressiveness_uplift"])]
        y += 0.69


def _add_recommendation_box(slide, df: pd.DataFrame) -> None:
    ranked = df[df["aggressiveness_uplift"] == "+60%"].sort_values("Budget_2026_EUR", ascending=False)
    top_markets = ranked["Market"].head(4).tolist()
    top_text = ", ".join(top_markets) if top_markets else "PCNA, PD, TRK, PKO"

    shape = slide.shapes.add_shape(1, Inches(10.05), Inches(4.86), Inches(2.9), Inches(1.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 118, 110)
    shape.line.color.rgb = RGBColor(15, 118, 110)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Recommendation"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = f"Prioritize aggressive funding in: {top_text}."
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = RGBColor(236, 253, 245)
    p2.font.name = "Calibri"

    p3 = tf.add_paragraph()
    p3.text = "Keep +40% markets on disciplined baseline until efficiency improves."
    p3.font.size = Pt(10)
    p3.font.color.rgb = RGBColor(236, 253, 245)
    p3.font.name = "Calibri"


def build_slide(df: pd.DataFrame) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.09))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bar.line.fill.background()

    _add_title(slide, "Recommendation: Market Clusters for 2026 Paid Media Aggressiveness")
    _add_subtitle(
        slide,
        "We assign +40% / +50% / +60% bands by market based on target pressure and scalability headroom.",
    )
    slide.shapes.add_picture(str(SCATTER_PNG), Inches(0.45), Inches(1.17), width=Inches(9.45), height=Inches(4.9))
    _add_cluster_summary(slide, df)
    _add_recommendation_box(slide, df)
    _add_method_box(slide)

    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.02), Inches(12.2), Inches(0.25))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = (
        "Note: CTG data is used as a directional proxy for paid-media responsiveness. "
        "Bands are guidance on urgency/importance, not deterministic allocation outputs."
    )
    p.font.size = Pt(9)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(107, 114, 128)
    p.alignment = PP_ALIGN.LEFT

    prs.save(OUTPUT_PPTX)


def main() -> None:
    df = load_data()
    build_scatter_image(df)
    build_slide(df)
    print(f"Saved slide: {OUTPUT_PPTX}")
    print(f"Saved chart image: {SCATTER_PNG}")


if __name__ == "__main__":
    main()
