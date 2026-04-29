from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from openpyxl import load_workbook
from matplotlib.lines import Line2D
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
WORKBOOK_PATH = BASE_DIR / "2026_03_03_RevFLight_Simplified.xlsx"
REGISTRATIONS_PATH = BASE_DIR / "Registrations over time split by Country.xlsx"
SCATTER_PNG = BASE_DIR / "reverse_funnel_aggressiveness_scatter.png"
OUTPUT_PPTX = BASE_DIR / "reverse_funnel_aggressiveness_slide.pptx"


def _extract_market_values(ws, row_idx: int) -> dict[str, float]:
    markets: dict[str, float] = {}
    for col in range(4, 17):  # D:P
        market = ws.cell(77, col).value
        value = ws.cell(row_idx, col).value
        if market in (None, "TOTAL"):
            continue
        if value is None:
            continue
        markets[str(market)] = float(value)
    return markets


def load_registrations_headroom() -> pd.DataFrame:
    country_map = {
        "PBR": "Brazil",
        "PCA": "Australia",
        "PCGB": "United Kingdom",
        "PCH": "Switzerland",
        "PCL": "Canada",
        "PCNA": "United States",
        "PD": "Germany",
        "PIB": ["Spain", "Portugal"],
        "PIB SPA": "Spain",
        "PIB POR": "Portugal",
        "PIT": "Italy",
        "PNO": "Norway",
        "POF": "France",
        "PPL": "Poland",
        "PTW": "Taiwan",
    }

    df = pd.read_excel(REGISTRATIONS_PATH, sheet_name="Sheet1", header=2)
    df = df[["Year", "Country", "Registrations"]].copy()
    df["Registrations"] = pd.to_numeric(df["Registrations"], errors="coerce")
    df = df.dropna(subset=["Year", "Country", "Registrations"])
    yearly = (
        df[df["Year"].isin([2024, 2025])]
        .groupby(["Country", "Year"], as_index=False)["Registrations"]
        .sum()
    )
    pivot = yearly.pivot(index="Country", columns="Year", values="Registrations").reset_index()
    pivot.columns.name = None
    pivot = pivot.rename(columns={2024: "registrations_2024", 2025: "registrations_2025"})
    for col in ["registrations_2024", "registrations_2025"]:
        if col not in pivot.columns:
            pivot[col] = 0.0

    rows = []
    for market, country in country_map.items():
        countries = country if isinstance(country, list) else [country]
        rec = pivot[pivot["Country"].isin(countries)]
        if rec.empty:
            r2024 = float("nan")
            r2025 = float("nan")
        else:
            r2024 = float(rec["registrations_2024"].sum(min_count=1))
            r2025 = float(rec["registrations_2025"].sum(min_count=1))
        growth = (r2025 / r2024 - 1.0) if pd.notna(r2024) and r2024 > 0 else float("nan")
        rows.append(
            {
                "Market": market,
                "Country": ",".join(countries),
                "registrations_2024": r2024,
                "registrations_2025": r2025,
                "registrations_yoy_growth": growth,
            }
        )
    return pd.DataFrame(rows)


def load_reverse_funnel_metrics() -> pd.DataFrame:
    wb = load_workbook(WORKBOOK_PATH, data_only=True)
    ws = wb["RFL"]

    order_intake = _extract_market_values(ws, 78)
    required_paid_dcfs = _extract_market_values(ws, 98)
    budget_2026 = _extract_market_values(ws, 110)

    records = []
    for market, budget in budget_2026.items():
        paid_dcfs = required_paid_dcfs.get(market)
        oi = order_intake.get(market)
        if paid_dcfs is None or oi is None or paid_dcfs <= 0:
            continue
        records.append(
            {
                "Market": market,
                "OrderIntake_2026": oi,
                "Budget_2026_EUR": budget,
                "RequiredPaidDCFS_2026": paid_dcfs,
                "CostPerRequiredPaidDCFS_EUR": budget / paid_dcfs,
                "PaidDCFSPereuro": paid_dcfs / budget if budget > 0 else 0.0,
            }
        )

    df = pd.DataFrame(records)
    reg = load_registrations_headroom()
    df = df.merge(reg, on="Market", how="left")

    # Headroom proxy from registrations:
    # - share gap: registrations share minus budget share (under-allocated demand if positive)
    # - momentum: YoY registrations growth as additional demand signal
    total_reg = df["registrations_2025"].sum(min_count=1)
    total_budget = df["Budget_2026_EUR"].sum(min_count=1)
    df["reg_share_2025"] = df["registrations_2025"] / total_reg if pd.notna(total_reg) and total_reg > 0 else float("nan")
    df["budget_share_2026"] = df["Budget_2026_EUR"] / total_budget if pd.notna(total_budget) and total_budget > 0 else float("nan")
    df["headroom_share_gap"] = df["reg_share_2025"] - df["budget_share_2026"]

    df["pressure_pct"] = df["RequiredPaidDCFS_2026"].rank(pct=True)
    df["efficiency_pct"] = df["PaidDCFSPereuro"].rank(pct=True)
    df["headroom_share_pct"] = df["headroom_share_gap"].rank(pct=True)
    df["headroom_growth_pct"] = df["registrations_yoy_growth"].rank(pct=True)
    df["headroom_pct"] = 0.7 * df["headroom_share_pct"] + 0.3 * df["headroom_growth_pct"]
    # Fill missing headroom with neutral midpoint so missing countries do not distort bands.
    df["headroom_pct"] = df["headroom_pct"].fillna(0.5)

    # Non-linear rule-based clustering (no weighted linear combination).
    df["pressure_band"] = pd.qcut(df["pressure_pct"], q=3, labels=["Low", "Mid", "High"]).astype(str)
    df["efficiency_band"] = pd.qcut(df["efficiency_pct"], q=3, labels=["Low", "Mid", "High"]).astype(str)
    df["headroom_band"] = pd.qcut(df["headroom_pct"], q=3, labels=["Low", "Mid", "High"]).astype(str)
    df["aggressiveness_uplift"] = df.apply(_assign_band_rule, axis=1)
    df["within_band_rank"] = df.apply(_within_band_rank, axis=1)
    return df.sort_values(["aggressiveness_uplift", "within_band_rank"], ascending=[False, False]).reset_index(drop=True)


def _assign_band_rule(row: pd.Series) -> str:
    pressure = row["pressure_band"]
    efficiency = row["efficiency_band"]
    headroom = row["headroom_band"]

    if pressure == "High" and (efficiency == "High" or headroom == "High"):
        return "+60%"
    if pressure == "Mid" and efficiency == "High" and headroom in {"High", "Mid"}:
        return "+60%"

    high_count = int(pressure == "High") + int(efficiency == "High") + int(headroom == "High")
    low_count = int(pressure == "Low") + int(efficiency == "Low") + int(headroom == "Low")
    if high_count >= 1 and low_count <= 1:
        return "+50%"
    if pressure == "Mid" and efficiency == "Mid":
        return "+50%"
    return "+40%"


def _within_band_rank(row: pd.Series) -> int:
    # Order inside each band without linear weighting:
    # pressure (most important) > efficiency > headroom.
    scale = {"Low": 0, "Mid": 1, "High": 2}
    return 100 * scale.get(row["pressure_band"], 0) + 10 * scale.get(row["efficiency_band"], 0) + scale.get(row["headroom_band"], 0)


def build_scatter(df: pd.DataFrame) -> None:
    colors = {"+60%": "#0F766E", "+50%": "#2563EB", "+40%": "#6B7280"}
    fig, ax = plt.subplots(figsize=(10.6, 5.6), dpi=210)
    fig.patch.set_facecolor("#FFFFFF")
    ax.set_facecolor("#F8FAFC")

    for band in ["+40%", "+50%", "+60%"]:
        sub = df[df["aggressiveness_uplift"] == band].copy()
        if sub.empty:
            continue
        # radius is driven by headroom percentile
        sizes = 180 + 900 * sub["headroom_pct"].fillna(0.5)
        ax.scatter(
            sub["pressure_pct"] * 100,
            sub["efficiency_pct"] * 100,
            s=sizes,
            c=colors[band],
            alpha=0.85,
            edgecolors="#FFFFFF",
            linewidths=1.2,
            label=band,
        )
        for _, row in sub.iterrows():
            ax.text(
                row["pressure_pct"] * 100 + 1.0,
                row["efficiency_pct"] * 100 + 0.8,
                str(row["Market"]),
                fontsize=8.3,
                color="#111827",
            )

    ax.grid(alpha=0.18, linewidth=0.8)
    ax.set_xlim(0, 102)
    ax.set_ylim(0, 102)
    ax.set_xlabel("Pressure percentile (higher = more required paid DCFS)", fontsize=9.5)
    ax.set_ylabel("Efficiency percentile (higher = more DCFS per EUR)", fontsize=9.5)
    ax.set_title("Market Clusters: Pressure vs Efficiency (Bubble size = Headroom)", fontsize=13, fontweight="bold", pad=10)
    legend_handles = [
        Line2D(
            [0],
            [0],
            marker="o",
            linestyle="None",
            markerfacecolor=colors["+40%"],
            markeredgecolor="#FFFFFF",
            markeredgewidth=0.8,
            markersize=9,
            label="+40%",
        ),
        Line2D(
            [0],
            [0],
            marker="o",
            linestyle="None",
            markerfacecolor=colors["+50%"],
            markeredgecolor="#FFFFFF",
            markeredgewidth=0.8,
            markersize=9,
            label="+50%",
        ),
        Line2D(
            [0],
            [0],
            marker="o",
            linestyle="None",
            markerfacecolor=colors["+60%"],
            markeredgecolor="#FFFFFF",
            markeredgewidth=0.8,
            markersize=9,
            label="+60%",
        ),
    ]
    ax.legend(
        handles=legend_handles,
        title="Uplift band",
        frameon=False,
        loc="lower right",
        fontsize=8.5,
        title_fontsize=9,
    )
    for spine in ax.spines.values():
        spine.set_color("#CBD5E1")

    plt.tight_layout(rect=[0.02, 0.04, 0.98, 0.95])
    fig.savefig(SCATTER_PNG, bbox_inches="tight")
    plt.close(fig)


def _add_header_bar(slide) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.09))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bar.line.fill.background()


def _add_title(slide) -> None:
    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.13), Inches(12.2), Inches(0.58))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Recommendation: 2026 Market Budget Aggressiveness (Reverse Funnel Only)"
    p.font.bold = True
    p.font.size = Pt(25)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(17, 24, 39)

    subtitle = slide.shapes.add_textbox(Inches(0.45), Inches(0.69), Inches(12.2), Inches(0.45))
    sf = subtitle.text_frame
    sf.clear()
    sp = sf.paragraphs[0]
    sp.text = "Method is based only on the RFL tab: required paid-media DCFS and budget-to-output efficiency by market."
    sp.font.size = Pt(12)
    sp.font.name = "Calibri"
    sp.font.color.rgb = RGBColor(55, 65, 81)


def _add_method_box(slide) -> None:
    shape = slide.shapes.add_shape(1, Inches(0.45), Inches(6.1), Inches(12.2), Inches(0.92))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(239, 246, 255)
    shape.line.color.rgb = RGBColor(209, 213, 219)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Method: (1) Delivery pressure = required paid-media DCFS from Reverse Funnel (RFL row 98). "
        "(2) Cost efficiency = 2026 budget / required paid-media DCFS (RFL row 110 over row 98). "
        "(3) Headroom = registrations share gap + registrations momentum. "
        "(4) Bands are assigned with rule-based clustering (no weighted linear score): +40%, +50%, +60%."
    )
    p.font.size = Pt(10.3)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(31, 41, 55)


def _add_band_summary(slide, df: pd.DataFrame) -> None:
    summary = (
        df.groupby("aggressiveness_uplift", dropna=False)
        .agg(
            markets=("Market", "count"),
            budget=("Budget_2026_EUR", "sum"),
            paid_dcfs=("RequiredPaidDCFS_2026", "sum"),
        )
        .reindex(["+60%", "+50%", "+40%"])
        .dropna(how="all")
        .reset_index()
    )
    color_map = {"+60%": RGBColor(15, 118, 110), "+50%": RGBColor(37, 99, 235), "+40%": RGBColor(107, 114, 128)}

    head = slide.shapes.add_textbox(Inches(9.95), Inches(1.2), Inches(2.95), Inches(0.4))
    htf = head.text_frame
    htf.clear()
    hp = htf.paragraphs[0]
    hp.text = "Uplift bands"
    hp.font.bold = True
    hp.font.size = Pt(12)
    hp.font.name = "Calibri"
    hp.font.color.rgb = RGBColor(17, 24, 39)

    y = 1.56
    for _, r in summary.iterrows():
        band = str(r["aggressiveness_uplift"])
        shape = slide.shapes.add_shape(1, Inches(9.95), Inches(y), Inches(2.95), Inches(0.72))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(229, 231, 235)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{band}  |  {int(r['markets'])} markets"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = color_map[band]
        p2 = tf.add_paragraph()
        p2.text = f"Budget: EUR {int(r['budget']):,} | Paid DCFS: {int(r['paid_dcfs']):,}"
        p2.font.size = Pt(9.5)
        p2.font.name = "Calibri"
        p2.font.color.rgb = RGBColor(75, 85, 99)
        y += 0.79


def _add_recommendation(slide, df: pd.DataFrame) -> None:
    top = df[df["aggressiveness_uplift"] == "+60%"]["Market"].head(5).tolist()
    top_text = ", ".join(top) if top else "n/a"

    shape = slide.shapes.add_shape(1, Inches(9.95), Inches(4.92), Inches(2.95), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 118, 110)
    shape.line.color.rgb = RGBColor(15, 118, 110)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Recommendation"
    p.font.bold = True
    p.font.size = Pt(12.5)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = f"Prioritize +60% markets first: {top_text}."
    p2.font.size = Pt(10.2)
    p2.font.name = "Calibri"
    p2.font.color.rgb = RGBColor(236, 253, 245)
    p3 = tf.add_paragraph()
    p3.text = "Use +50% as controlled scale and +40% as baseline optimization."
    p3.font.size = Pt(9.5)
    p3.font.name = "Calibri"
    p3.font.color.rgb = RGBColor(236, 253, 245)


def build_slide(df: pd.DataFrame) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_header_bar(slide)
    _add_title(slide)
    slide.shapes.add_picture(str(SCATTER_PNG), Inches(0.45), Inches(1.18), width=Inches(9.4), height=Inches(4.86))
    _add_band_summary(slide, df)
    _add_recommendation(slide, df)
    _add_method_box(slide)

    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.02), Inches(12.2), Inches(0.24))
    ft = footer.text_frame
    ft.clear()
    p = ft.paragraphs[0]
    p.text = "Data source: 2026_03_03_RevFLight_Simplified.xlsx (RFL tab). This cut is intentionally reverse-funnel-only."
    p.font.size = Pt(9)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(107, 114, 128)
    p.alignment = PP_ALIGN.LEFT

    prs.save(OUTPUT_PPTX)


def main() -> None:
    df = load_reverse_funnel_metrics()
    build_scatter(df)
    build_slide(df)
    print(f"Saved chart image: {SCATTER_PNG}")
    print(f"Saved slide: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
