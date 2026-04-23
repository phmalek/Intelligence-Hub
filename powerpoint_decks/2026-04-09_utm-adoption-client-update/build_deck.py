from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_PATH = BASE_DIR / "utm_adoption_client_update.pptx"


BG = RGBColor(18, 20, 24)
SURFACE = RGBColor(28, 31, 36)
SURFACE_2 = RGBColor(36, 40, 46)
TEXT = RGBColor(244, 245, 247)
MUTED = RGBColor(165, 171, 180)
LINE = RGBColor(72, 78, 88)
ACCENT = RGBColor(86, 191, 255)
WARM = RGBColor(244, 175, 79)
GREEN = RGBColor(113, 201, 153)
RED = RGBColor(227, 101, 101)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_box(slide, left, top, width, height, fill, line=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_text(slide, left, top, width, height, text, size, color=TEXT, bold=False, font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullet_list(slide, left, top, width, items, size=14, color=TEXT, bullet_color=None, line_gap=0.3):
    box = slide.shapes.add_textbox(left, top, width, Inches(max(0.3, 0.26 * len(items))))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.bullet = True
    return box


def add_chip(slide, left, top, width, text, fill, text_color=TEXT, line=None):
    add_box(slide, left, top, width, Inches(0.28), fill, line=line, radius=True)
    add_text(slide, left + Inches(0.08), top + Inches(0.035), width - Inches(0.16), Inches(0.2), text, 10, text_color, bold=True)


def add_divider(slide, y):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), y, Inches(12.2), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def footer(slide):
    add_text(slide, Inches(4.4), Inches(7.0), Inches(4.5), Inches(0.18), "Confidential - client update on UTM adoption follow-up", 8, MUTED, False)


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(6.8), Inches(0.3), "UTM ADOPTION FOLLOW-UP", 14, ACCENT, True)
    add_text(slide, Inches(0.7), Inches(0.95), Inches(7.2), Inches(1.0), "Current market status,\nresponse patterns and next steps", 28, TEXT, True)
    add_text(slide, Inches(0.72), Inches(2.1), Inches(6.0), Inches(0.4), "Built from the latest global follow-up tracker", 14, MUTED)
    add_divider(slide, Inches(2.65))
    panel = add_box(slide, Inches(8.05), Inches(0.75), Inches(4.3), Inches(4.8), SURFACE, line=LINE)
    add_text(slide, Inches(8.35), Inches(1.05), Inches(2.0), Inches(0.25), "IN THIS UPDATE", 12, MUTED, True)
    add_text(slide, Inches(8.35), Inches(1.5), Inches(3.5), Inches(0.4), "Where follow-up stands", 20, TEXT, True)
    add_text(slide, Inches(8.35), Inches(2.08), Inches(3.5), Inches(0.4), "How markets responded", 20, TEXT, True)
    add_text(slide, Inches(8.35), Inches(2.66), Inches(3.5), Inches(0.4), "What is blocking closure", 20, TEXT, True)
    add_text(slide, Inches(8.35), Inches(3.24), Inches(3.5), Inches(0.4), "What happens next", 20, TEXT, True)
    add_chip(slide, Inches(0.72), Inches(6.25), Inches(2.45), "BLUE = INTERNAL COORDINATION", ACCENT, BG)
    footer(slide)


def kpi_card(slide, left, top, big, label, fill):
    add_box(slide, left, top, Inches(2.65), Inches(1.2), fill, line=LINE)
    add_text(slide, left + Inches(0.2), top + Inches(0.18), Inches(1.1), Inches(0.4), str(big), 28, TEXT, True)
    add_text(slide, left + Inches(0.2), top + Inches(0.68), Inches(2.1), Inches(0.25), label, 11, MUTED, True)


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.7), Inches(0.45), Inches(5.0), Inches(0.35), "WHERE WE STAND", 24, TEXT, True)
    kpi_card(slide, Inches(0.7), Inches(1.0), "11", "Markets tracked", SURFACE)
    kpi_card(slide, Inches(3.55), Inches(1.0), "7", "Engaged and fixing", SURFACE)
    kpi_card(slide, Inches(6.4), Inches(1.0), "2", "Awaiting response", SURFACE)
    kpi_card(slide, Inches(9.25), Inches(1.0), "1", "No active scope", SURFACE)

    add_box(slide, Inches(0.7), Inches(2.55), Inches(12.0), Inches(3.8), SURFACE, line=LINE)
    blocks = [
        ("ENGAGED / FIXING", GREEN, ["Australia (PCA)", "UK (PCGB)", "Switzerland (PCH)", "Poland (PPL)", "Norway (PNO)", "PCEE", "LATAM (CL)"], Inches(1.0)),
        ("NEEDS EVIDENCE / ALIGNMENT", WARM, ["France (POF)"], Inches(4.25)),
        ("AWAITING RESPONSE", RED, ["Portugal", "South Korea"], Inches(6.55)),
        ("NO ACTIVE SCOPE", MUTED, ["MENA (PME)"], Inches(9.35)),
    ]
    for title, color, items, x in blocks:
        add_text(slide, x, Inches(2.9), Inches(2.2), Inches(0.22), title, 10, color, True)
        for i, item in enumerate(items):
            add_box(slide, x, Inches(3.25) + Inches(0.38 * i), Inches(2.05), Inches(0.28), SURFACE_2, line=LINE)
            add_text(slide, x + Inches(0.15), Inches(3.31) + Inches(0.38 * i), Inches(1.9), Inches(0.2), item, 11, TEXT)

    add_chip(slide, Inches(9.6), Inches(5.55), Inches(2.1), "PCGB - builder support", ACCENT, BG)
    add_chip(slide, Inches(9.6), Inches(5.9), Inches(1.95), "PNO - Planit onboarding", ACCENT, BG)
    add_chip(slide, Inches(9.6), Inches(6.25), Inches(1.85), "PCEE - Artbot remap", ACCENT, BG)
    footer(slide)


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.7), Inches(0.45), Inches(6.0), Inches(0.35), "HOW MARKETS RESPONDED", 24, TEXT, True)
    cols = [
        ("ACKNOWLEDGED AND ENGAGED", "Markets accepted the issue and are now working through fixes or validation.", ["Australia", "UK", "Switzerland", "Poland", "LATAM"], GREEN),
        ("ASKING FOR PROOF OR CLARITY", "Some markets challenged the diagnosis and now need evidence-based closure.", ["France", "Australia", "Switzerland"], WARM),
        ("NEEDS STRUCTURE OR ONBOARDING", "In some cases the blocker is setup maturity, not willingness to adopt.", ["Norway", "PCEE", "MENA"], ACCENT),
    ]
    lefts = [Inches(0.7), Inches(4.45), Inches(8.2)]
    for (title, body, items, color), left in zip(cols, lefts):
        add_box(slide, left, Inches(1.2), Inches(3.1), Inches(4.7), SURFACE, line=LINE)
        add_text(slide, left + Inches(0.22), Inches(1.48), Inches(2.7), Inches(0.35), title, 12, color, True)
        add_text(slide, left + Inches(0.22), Inches(1.95), Inches(2.6), Inches(0.8), body, 14, TEXT, False)
        for i, item in enumerate(items):
            add_chip(slide, left + Inches(0.22), Inches(3.15) + Inches(0.42 * i), Inches(1.7), item, SURFACE_2, TEXT, line=LINE)
    add_box(slide, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.7), ACCENT, line=None)
    add_text(slide, Inches(0.95), Inches(6.42), Inches(11.2), Inches(0.25), "Several open items now depend on internal coordination across Planit, builder usage, trafficking logic or placement-ID structure.", 13, BG, True)
    footer(slide)


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.7), Inches(0.45), Inches(6.0), Inches(0.35), "WHAT IS BLOCKING CLOSURE", 24, TEXT, True)
    cards = [
        ("TAGGING GAPS", "Live campaigns or placements are still not consistently visible in GA4."),
        ("PLANIT / EXECUTION MISALIGNMENT", "Planned activity and live activity do not fully reconcile."),
        ("STRUCTURE ISSUES", "Placement IDs, package-level tagging or repeated IDs break clean visibility."),
        ("TRACKING-LOGIC MISUNDERSTANDING", "Some markets need alignment on expected source logic across GA4, 1R and CM360."),
    ]
    positions = [(Inches(0.7), Inches(1.2)), (Inches(6.5), Inches(1.2)), (Inches(0.7), Inches(3.8)), (Inches(6.5), Inches(3.8))]
    for (title, body), (left, top) in zip(cards, positions):
        add_box(slide, left, top, Inches(5.4), Inches(2.2), SURFACE, line=LINE)
        add_text(slide, left + Inches(0.24), top + Inches(0.24), Inches(3.7), Inches(0.3), title, 12, ACCENT, True)
        add_text(slide, left + Inches(0.24), top + Inches(0.78), Inches(4.5), Inches(0.9), body, 16, TEXT)
    add_text(slide, Inches(0.7), Inches(6.55), Inches(11.0), Inches(0.3), "The pattern is repeatable: this is now less about awareness and more about implementation discipline.", 16, MUTED, True)
    footer(slide)


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.7), Inches(0.45), Inches(6.0), Inches(0.35), "NEXT STEPS TO CLOSE THE LOOP", 24, TEXT, True)
    add_box(slide, Inches(0.7), Inches(1.2), Inches(5.7), Inches(4.9), SURFACE, line=ACCENT)
    add_box(slide, Inches(6.6), Inches(1.2), Inches(5.7), Inches(4.9), SURFACE, line=LINE)
    add_text(slide, Inches(0.95), Inches(1.45), Inches(2.0), Inches(0.25), "GLOBAL TEAM", 13, ACCENT, True)
    add_bullet_list(slide, Inches(0.95), Inches(1.95), Inches(4.7), [
        "Provide market-specific evidence where diagnosis is disputed",
        "Support builder, Planit and trafficking setup where structure is missing",
        "Clarify expected tracking logic across GA4, 1R and CM360",
        "Validate fixes once markets confirm changes",
    ], size=15)
    add_text(slide, Inches(6.85), Inches(1.45), Inches(2.4), Inches(0.25), "LOCAL MARKETS", 13, TEXT, True)
    add_bullet_list(slide, Inches(6.85), Inches(1.95), Inches(4.7), [
        "Correct remaining UTM gaps",
        "Align Planit with live activity",
        "Resolve placement-ID and setup issues",
        "Confirm closure with evidence, not assumption",
    ], size=15)
    add_box(slide, Inches(0.7), Inches(6.3), Inches(11.6), Inches(0.72), SURFACE_2, line=LINE)
    add_text(slide, Inches(0.95), Inches(6.53), Inches(10.8), Inches(0.2), "Immediate focus: convert open discussions into validated closure by market type.", 15, TEXT, True)
    footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    prs.save(OUTPUT_PATH)
    print(f'Wrote {OUTPUT_PATH}')


if __name__ == "__main__":
    main()
