from __future__ import annotations

import argparse
import copy
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement


@dataclass
class CellTarget:
    row_idx: int
    col_idx: int
    market: str
    platform: str
    center: Tuple[float, float]


def _norm(text: object) -> str:
    if text is None:
        return ""
    return re.sub(r"\s+", " ", str(text).strip()).lower()


def _platform_key(text: object) -> str:
    s = _norm(text)
    s = s.replace("&", "and")
    aliases = {
        "meta": "facebook",
        "facebook ads": "facebook",
        "google": "google ads",
        "googlead": "google ads",
        "googleads": "google ads",
        "linkedin": "linkedin",
        "linked in": "linkedin",
        "tiktok": "tiktok",
        "tik tok": "tiktok",
        "dv 360": "dv360",
        "direct buys": "direct buys",
    }
    squashed = re.sub(r"[^a-z0-9 ]+", "", s).strip()
    squashed_nospace = squashed.replace(" ", "")
    if squashed in aliases:
        return aliases[squashed]
    if squashed_nospace in aliases:
        return aliases[squashed_nospace]
    return squashed


def _duplicate_slide(prs: Presentation, slide_index: int) -> int:
    """Duplicate slide by copying shape XML into a new slide."""
    source = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6]  # blank
    new_slide = prs.slides.add_slide(blank_layout)

    # Copy all shape elements
    for shape in source.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Copy known relationships (images, ole, etc.)
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)
        except Exception:
            # best effort; XML copy usually preserves enough for rendering
            pass

    return len(prs.slides) - 1


def _get_table(slide):
    for sh in slide.shapes:
        if sh.has_table:
            return sh
    raise RuntimeError("No table found on target slide.")


def _cell_centers(table_shape) -> List[CellTarget]:
    tbl = table_shape.table
    left = table_shape.left
    top = table_shape.top

    col_widths = [c.width for c in tbl.columns]
    row_heights = [r.height for r in tbl.rows]

    # cumulative positions
    x_edges = [left]
    for w in col_widths:
        x_edges.append(x_edges[-1] + w)
    y_edges = [top]
    for h in row_heights:
        y_edges.append(y_edges[-1] + h)

    targets: List[CellTarget] = []
    for r in range(1, len(tbl.rows)):  # skip header row
        market = tbl.cell(r, 0).text.strip()
        if not market:
            continue
        for c in range(1, len(tbl.columns)):  # skip market column
            platform = tbl.cell(0, c).text.strip()
            if not platform:
                continue
            cx = (x_edges[c] + x_edges[c + 1]) / 2
            cy = (y_edges[r] + y_edges[r + 1]) / 2
            targets.append(CellTarget(r, c, market, platform, (cx, cy)))
    return targets


def _ellipse_centers(slide, table_shape):
    t_left = table_shape.left
    t_top = table_shape.top
    t_right = t_left + table_shape.width
    t_bottom = t_top + table_shape.height

    circles = []
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if "Ellipse" not in sh.name:
            continue
        cx = sh.left + sh.width / 2
        cy = sh.top + sh.height / 2
        if t_left <= cx <= t_right and t_top <= cy <= t_bottom:
            circles.append((sh, (cx, cy)))
    return circles


def _assign_circles_to_cells(
    targets: List[CellTarget], circles: List[Tuple[object, Tuple[float, float]]]
) -> Dict[Tuple[str, str], object]:
    """
    Greedy nearest-neighbor one-to-one assignment.
    Returns key (market, platform) -> shape.
    """
    unassigned = circles[:]
    mapping: Dict[Tuple[str, str], object] = {}
    # predictable order: row then col
    targets_sorted = sorted(targets, key=lambda t: (t.row_idx, t.col_idx))

    for t in targets_sorted:
        if not unassigned:
            break
        best_i = None
        best_d = None
        tx, ty = t.center
        for i, (_, (cx, cy)) in enumerate(unassigned):
            d = (cx - tx) ** 2 + (cy - ty) ** 2
            if best_d is None or d < best_d:
                best_d = d
                best_i = i
        sh, _ = unassigned.pop(best_i)  # one circle per cell
        mapping[(_norm(t.market).upper(), _platform_key(t.platform))] = sh
    return mapping


def _status_to_rgb(status: str) -> Optional[RGBColor]:
    s = _norm(status)
    if not s:
        return None

    # response traffic light values
    if "green" in s or "full implementation" in s or "full" == s:
        return RGBColor(0x92, 0xD0, 0x50)
    if "amber" in s or "partial" in s or "check still in progress" in s:
        return RGBColor(0xFF, 0xBE, 0x00)
    if "red" in s or "other concept" in s:
        return RGBColor(0xFF, 0x00, 0x00)
    if "not live" in s or "no live" in s or "no active" in s or "hypothesis" in s:
        return RGBColor(0x00, 0x00, 0x00)
    if "unknown" in s:
        return RGBColor(0xFF, 0xBE, 0x00)
    return None


def _load_status_map(consolidated_path: Path) -> Dict[Tuple[str, str], str]:
    df = pd.read_excel(consolidated_path) if consolidated_path.suffix.lower() == ".xlsx" else pd.read_csv(consolidated_path)
    required = {"Market", "Platform (from input)"}
    if not required.issubset(df.columns):
        raise RuntimeError(f"Consolidated data missing required columns: {required}")

    status_map: Dict[Tuple[str, str], str] = {}
    for _, row in df.iterrows():
        market = str(row.get("Market", "")).strip().upper()
        platform = _platform_key(row.get("Platform (from input)", ""))
        if not market or not platform:
            continue
        # Strict rule: update only when response traffic light is provided.
        # If response is missing, leave existing deck color unchanged.
        response_tl = row.get("Response Traffic Light")
        if pd.isna(response_tl) or not str(response_tl).strip():
            continue
        status = response_tl
        status_map[(market, platform)] = str(status)
    return status_map


def _set_fill(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.width = 0


def update_trafficlights(
    pptx_path: Path,
    consolidated_path: Path,
    output_path: Path,
    source_slide_index_1_based: int = 2,
    duplicate: bool = True,
) -> None:
    prs = Presentation(pptx_path)
    src_idx = source_slide_index_1_based - 1
    if duplicate:
        target_idx = _duplicate_slide(prs, src_idx)
    else:
        target_idx = src_idx
    slide = prs.slides[target_idx]

    table_shape = _get_table(slide)
    targets = _cell_centers(table_shape)
    circles = _ellipse_centers(slide, table_shape)
    shape_map = _assign_circles_to_cells(targets, circles)
    status_map = _load_status_map(consolidated_path)

    updated = 0
    unmatched = 0
    for key, status in status_map.items():
        # key is (MARKET, platform_norm)
        sh = shape_map.get((key[0], key[1]))
        if sh is None:
            unmatched += 1
            continue
        rgb = _status_to_rgb(status)
        if rgb is None:
            continue
        _set_fill(sh, rgb)
        updated += 1

    prs.save(output_path)
    if duplicate:
        print(f"Duplicated slide index: {source_slide_index_1_based} -> new slide index: {target_idx + 1}")
    else:
        print(f"Updated slide index in place: {source_slide_index_1_based}")
    print(f"Traffic lights updated: {updated}")
    print(f"Status rows unmatched to grid: {unmatched}")
    print(f"Saved: {output_path}")


def main():
    ap = argparse.ArgumentParser(
        description="Duplicate traffic-light slide and only update existing circle fill colors."
    )
    ap.add_argument("--pptx", required=True, help="Input pptx path")
    ap.add_argument("--consolidated", required=True, help="Consolidated response file (csv or xlsx)")
    ap.add_argument("--output", required=True, help="Output pptx path")
    ap.add_argument("--slide-index", type=int, default=2, help="1-based index of source slide to duplicate (default: 2)")
    ap.add_argument("--in-place", action="store_true", help="Update the target slide in place (no duplication)")
    args = ap.parse_args()

    update_trafficlights(
        pptx_path=Path(args.pptx),
        consolidated_path=Path(args.consolidated),
        output_path=Path(args.output),
        source_slide_index_1_based=args.slide_index,
        duplicate=not args.in_place,
    )


if __name__ == "__main__":
    main()
